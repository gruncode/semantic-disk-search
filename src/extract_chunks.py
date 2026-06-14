#!/usr/bin/env python3
"""Extract text chunks from documents for GPU embedding.
Runs locally — outputs JSON Lines file with chunk text + metadata."""
import hashlib
import json
import os
import re
import subprocess
import sys
import zipfile
from pathlib import Path

CHUNK_SIZE = 500
CHUNK_OVERLAP = 100
MAX_TEXT_CHARS = 30000
MIN_CHUNK_LEN = 30

# OCR sidecar extensions in quality order (best → fallback).
# Used both for extraction priority and for skip-list during file scan.
OCR_SIDECARS = [
    '.ocr.2.5-pro.md',
    '.ocr.gemini-pro.md',
    '.ocr.gemini-pro.txt',
    '.ocr.docai.fmt.md',
    '.ocr.docai.md',
    '.ocr.gemini.md',
    '.ocr.gemini.txt',
    '.ocr.txt',
]
# All sidecar suffixes to skip during scan (so they're not indexed as standalone docs)
_OCR_SKIP = tuple(OCR_SIDECARS + [
    '.ocr.test-formparser.txt', '.ocr.test-docai.txt', '.ocr.retest-pro.txt',
    '.ocr.xlsx', '.ocr.docx',
])

TEXT_EXTS = {'.txt', '.md', '.rtf', '.html', '.htm', '.csv', '.log'}
OFFICE_EXTS = {'.odt', '.ods', '.docx', '.xlsx'}
IMAGE_EXTS = {'.jpg', '.jpeg', '.png', '.tif', '.tiff', '.bmp'}
XLS_EXTS = {'.xls'}
PPTX_EXTS = {'.pptx'}
ALL_EXTS = TEXT_EXTS | OFFICE_EXTS | IMAGE_EXTS | XLS_EXTS | PPTX_EXTS | {'.pdf', '.doc', '.ppt'}


def extract_text(path):
    p = Path(path)
    if not p.exists() or not p.is_file():
        return ''
    ext = p.suffix.lower()
    try:
        if ext in TEXT_EXTS:
            return p.read_text(errors='ignore')[:MAX_TEXT_CHARS]
        if ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            text = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) if c is not None else '' for c in row]
                    if any(v for v in vals):
                        text.append('\t'.join(vals))
            wb.close()
            return '\n'.join(text)[:MAX_TEXT_CHARS]
        if ext in XLS_EXTS:
            import xlrd
            wb = xlrd.open_workbook(path)
            text = []
            for ws in wb.sheets():
                for row_idx in range(ws.nrows):
                    vals = [str(ws.cell_value(row_idx, c)) for c in range(ws.ncols)]
                    if any(v for v in vals):
                        text.append('\t'.join(vals))
            return '\n'.join(text)[:MAX_TEXT_CHARS]
        if ext in PPTX_EXTS:
            from pptx import Presentation
            prs = Presentation(path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                text.append(t)
            return '\n'.join(text)[:MAX_TEXT_CHARS]
        if ext in OFFICE_EXTS:
            z = zipfile.ZipFile(path)
            txt = []
            for n in z.namelist():
                if n.endswith(('content.xml', 'document.xml', 'sharedStrings.xml')):
                    txt.append(z.read(n).decode('utf-8', 'ignore'))
            return re.sub(r'<[^>]+>', ' ', ' '.join(txt))[:MAX_TEXT_CHARS]
        if ext == '.pdf':
            for suffix in OCR_SIDECARS:
                ocr_cache = path + suffix
                if os.path.exists(ocr_cache):
                    text = Path(ocr_cache).read_text(errors='ignore').strip()
                    if len(text) >= MIN_CHUNK_LEN:
                        return text[:MAX_TEXT_CHARS]
            cp = subprocess.run(['pdftotext', '-q', '-nopgbrk', path, '-'],
                                capture_output=True, timeout=30)
            text = cp.stdout.decode('utf-8', errors='replace').strip()
            if len(text) >= MIN_CHUNK_LEN:
                return text[:MAX_TEXT_CHARS]
            return ''
        if ext == '.doc':
            cp = subprocess.run(['antiword', path], capture_output=True, timeout=15)
            return cp.stdout.decode('utf-8', errors='replace')[:MAX_TEXT_CHARS]
        if ext in IMAGE_EXTS:
            for suffix in OCR_SIDECARS:
                ocr_cache = path + suffix
                if os.path.exists(ocr_cache):
                    text = Path(ocr_cache).read_text(errors='ignore').strip()
                    if len(text) >= MIN_CHUNK_LEN:
                        return text[:MAX_TEXT_CHARS]
            return ''
    except Exception:
        pass
    return ''


def chunk_text(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP, section_aware=False):
    text = text.strip()
    if section_aware and re.search(r'\n#{1,3} ', text):
        sections = re.split(r'\n(?=#{1,3} )', text)
        chunks = []
        for s in sections:
            s = s.strip()
            if s:
                chunks.extend(chunk_text(s, chunk_size, overlap, section_aware=False))
        return chunks or chunk_text(text, chunk_size, overlap, section_aware=False)
    if len(text) <= chunk_size:
        return [text] if len(text) >= MIN_CHUNK_LEN else []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if len(chunk) >= MIN_CHUNK_LEN:
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks


def main():
    conf_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.expanduser('~/.recoll')
    output_file = sys.argv[2] if len(sys.argv) > 2 else '/tmp/chunks.jsonl'

    # Read topdirs from recoll.conf
    topdirs = []
    conf_file = os.path.join(conf_dir, 'recoll.conf')
    if os.path.exists(conf_file):
        with open(conf_file) as f:
            for line in f:
                if line.strip().startswith('topdirs'):
                    topdirs = line.split('=', 1)[1].strip().split()
                    break
    if not topdirs:
        topdirs = [os.path.expanduser('~/DATA')]

    # Scan files
    paths = []
    for topdir in topdirs:
        for root, dirs, files in os.walk(topdir, followlinks=True):
            dirs[:] = [d for d in dirs if d not in
                       {'.git', '.svn', 'node_modules', '__pycache__',
                        '.cache', '.thumbnails', '.Trash-0', '.Trash-1000'}]
            for fname in files:
                if fname.lower().endswith(_OCR_SKIP):
                    continue
                ext = Path(fname).suffix.lower()
                if ext in ALL_EXTS:
                    full = os.path.join(root, fname)
                    try:
                        if os.path.getsize(full) > 100:
                            paths.append(full)
                    except OSError:
                        pass

    print(f'Found {len(paths)} documents in {topdirs}', file=sys.stderr)

    total_chunks = 0
    skipped = 0
    with open(output_file, 'w') as out:
        for i, path in enumerate(paths):
            if i % 200 == 0:
                print(f'\r  Extracting: {i}/{len(paths)} ({total_chunks} chunks, {skipped} skipped)...',
                      end='', file=sys.stderr)

            text = extract_text(path)
            text = text.encode('utf-8', errors='ignore').decode('utf-8').strip()
            if len(text) < MIN_CHUNK_LEN:
                skipped += 1
                continue

            ext = Path(path).suffix.lower()
            # Get path relative to topdir for folder tag and filepath prefix
            rel = path
            folder = ''
            for td in topdirs:
                if path.startswith(td):
                    rel = path[len(td):].lstrip('/')
                    parts = rel.split('/')
                    folder = parts[0] if parts else ''
                    break
            # Short path tag: last 4 components (keeps context, avoids very long prefixes)
            rel_parts = rel.strip('/').split('/')
            path_tag = '/'.join(rel_parts[-4:]) if len(rel_parts) >= 4 else rel

            # Section-aware chunking for markdown content (OCR .md sidecars, .md files)
            has_sections = bool(re.search(r'\n#{1,3} ', text))
            chunks = chunk_text(text, section_aware=has_sections)
            if not chunks:
                skipped += 1
                continue

            mtime = str(int(os.path.getmtime(path))) if os.path.exists(path) else '0'
            chunk_id_base = hashlib.md5(path.encode('utf-8', errors='ignore')).hexdigest()[:12]
            # Sanitize path strings: Linux filenames can contain non-UTF-8 bytes which Python
            # represents as lone surrogates (\udcXX) via surrogateescaped — not JSON-serializable.
            path_safe     = path.encode('utf-8', errors='ignore').decode('utf-8')
            path_tag_safe = path_tag.encode('utf-8', errors='ignore').decode('utf-8')
            folder_safe   = folder.encode('utf-8', errors='ignore').decode('utf-8')

            for ci, chunk in enumerate(chunks):
                record = {
                    'id': f'{chunk_id_base}_{ci}',
                    'text': f'[{path_tag_safe}]\n{chunk}',
                    'path': path_safe,
                    'folder': folder_safe,
                    'ext': ext,
                    'mtime': mtime,
                    'chunk_idx': ci,
                    'total_chunks': len(chunks),
                }
                out.write(json.dumps(record, ensure_ascii=False) + '\n')
                total_chunks += 1

    print(f'\nDone: {total_chunks} chunks from {len(paths) - skipped} docs → {output_file}', file=sys.stderr)
    print(f'File size: {os.path.getsize(output_file) / 1024 / 1024:.1f} MB', file=sys.stderr)


if __name__ == '__main__':
    main()
